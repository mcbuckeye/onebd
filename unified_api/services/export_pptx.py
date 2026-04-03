"""Export comp sets and DD packages to PowerPoint and PDF."""
import io
import structlog
from typing import List, Dict, Any, Optional

logger = structlog.get_logger(__name__)


def generate_comp_pptx(title: str, criteria: Dict, deals: List[Dict], stats: Dict) -> bytes:
    """Generate a PowerPoint presentation for a comp set."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"Comparable Deal Analysis"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
    p2.alignment = PP_ALIGN.CENTER
    
    # Stats slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary Statistics"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x60, 0xa5, 0xfa)
    p.font.bold = True
    
    stat_items = [
        f"Total Comparable Deals: {stats.get('count', 0)}",
        f"Disclosed Values: {stats.get('disclosed', 0)} ({stats.get('disclosure_rate', 'N/A')}%)",
        f"Median Value: ${stats.get('median', 0):,.0f}M" if stats.get('median') else "Median Value: N/A",
        f"Mean Value: ${stats.get('mean', 0):,.0f}M" if stats.get('mean') else "Mean Value: N/A",
        f"Range: ${stats.get('min', 0):,.0f}M – ${stats.get('max', 0):,.0f}M" if stats.get('min') else "Range: N/A",
    ]
    
    for i, item in enumerate(stat_items):
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5 + i * 0.8), Inches(10), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xcb, 0xd5, 0xe1)
    
    # Deals table slide(s) - 8 deals per slide
    for chunk_start in range(0, len(deals), 8):
        chunk = deals[chunk_start:chunk_start + 8]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Comparable Deals ({chunk_start + 1}-{min(chunk_start + 8, len(deals))} of {len(deals)})"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x60, 0xa5, 0xfa)
        
        # Table
        rows = len(chunk) + 1  # header + data
        cols = 5
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1), Inches(12), Inches(rows * 0.6))
        table = table_shape.table
        
        headers = ['Deal', 'Principal → Partner', 'Phase', 'Value ($M)', 'Date']
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
                paragraph.font.bold = True
        
        for i, deal in enumerate(chunk):
            row_idx = i + 1
            values = [
                deal.get('title', '')[:60],
                f"{deal.get('principal_company', '?')} → {deal.get('partner_company', '?')}",
                deal.get('phase', '—'),
                f"${deal.get('total_value', 0):,.0f}" if deal.get('total_value') else '—',
                deal.get('date_start', '—'),
            ]
            for j, val in enumerate(values):
                cell = table.cell(row_idx, j)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(9)
                    paragraph.font.color.rgb = RGBColor(0xcb, 0xd5, 0xe1)
    
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_dd_pdf(company_name: str, dd_package: Dict) -> bytes:
    """Generate a PDF for a DD package."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.units import inch
    
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, leftMargin=inch, rightMargin=inch)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle('DDTitle', parent=styles['Title'], fontSize=24, textColor=HexColor('#1e293b'))
    heading_style = ParagraphStyle('DDHeading', parent=styles['Heading2'], fontSize=14, textColor=HexColor('#3b82f6'), spaceAfter=12)
    body_style = ParagraphStyle('DDBody', parent=styles['Normal'], fontSize=10, textColor=HexColor('#334155'), leading=14)
    
    story = []
    
    # Title
    story.append(Paragraph(f"Due Diligence Report: {company_name}", title_style))
    story.append(Spacer(1, 24))
    
    # Risk flags
    risk_flags = dd_package.get('risk_flags', [])
    if risk_flags:
        story.append(Paragraph("Risk Assessment", heading_style))
        for flag in risk_flags:
            severity = flag.get('severity', 'info')
            prefix = {'high': '🔴', 'medium': '🟡', 'low': '🟢'}.get(severity, 'ℹ️')
            story.append(Paragraph(f"{prefix} {flag.get('flag', '')}", body_style))
        story.append(Spacer(1, 16))
    
    # Sections
    for section in dd_package.get('sections', []):
        story.append(Paragraph(section.get('title', ''), heading_style))
        
        content = section.get('content')
        if isinstance(content, str):
            story.append(Paragraph(content, body_style))
        elif isinstance(content, dict):
            for key, val in content.items():
                if key != 'id':
                    story.append(Paragraph(f"<b>{key.replace('_', ' ').title()}:</b> {val or '—'}", body_style))
        elif isinstance(content, list):
            for item in content[:20]:
                if isinstance(item, dict):
                    text = item.get('title') or item.get('name') or str(item)
                    story.append(Paragraph(f"• {text}", body_style))
        
        story.append(Spacer(1, 12))
    
    doc.build(story)
    return buf.getvalue()
